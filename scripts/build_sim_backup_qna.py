#!/usr/bin/env python3
"""Build Deniz Ege's personal SIMULATION/methodology jury-Q&A backup deck.

Deliberately SIMPLE and PLAIN (white background, navy title band, no 3D, no glass
effects) so it contrasts with the animated main deck and prints cleanly. Scoped to the
simulation/methodology section. Every number is taken from the verified output-file
fact sheet (output/policy_stats.json, validation_report.json, zwm92_summary.json,
timing_study_f400.json, preprocess_stats.json, src/config.py).

Output: defense/FENG498_SimBackup_DenizEge.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "defense" / "FENG498_SimBackup_DenizEge.pptx"

NAVY = RGBColor(0x0F, 0x27, 0x47)
GREEN = RGBColor(0x21, 0x8A, 0x3A)
INK = RGBColor(0x22, 0x2A, 0x33)
GREY = RGBColor(0x5B, 0x66, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def title_slide():
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, SW, SH, NAVY)
    tf = _box(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.2))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Methodology & Simulation"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = "Backup Q&A, anticipated jury questions"
    r.font.size = Pt(22); r.font.color.rgb = RGBColor(0x9F, 0xC9, 0xA8)
    p3 = tf.add_paragraph(); p3.space_before = Pt(18)
    r = p3.add_run(); r.text = "Deniz Ege Memetoglu  |  FENG 498  |  Section 3"
    r.font.size = Pt(16); r.font.color.rgb = RGBColor(0xC7, 0xD2, 0xDD)
    f = _box(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
    rr = f.paragraphs[0].add_run()
    rr.text = "Plain backup sheet. Open in a second tab during Q&A. Numbers verified from the simulation output files."
    rr.font.size = Pt(11); rr.font.color.rgb = RGBColor(0x8A, 0x9A, 0xAB); rr.font.italic = True


def qna_slide(num, question, answer, bullets):
    s = prs.slides.add_slide(BLANK)
    # title band
    band = _rect(s, 0, 0, SW, Inches(1.5), NAVY)
    tf = band.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6); tf.margin_right = Inches(0.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"Q{num}.  {question}"
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    # short answer
    atf = _box(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(0.9))
    ap = atf.paragraphs[0]
    r = ap.add_run(); r.text = answer
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = GREEN
    # bullets
    btf = _box(s, Inches(0.7), Inches(2.75), Inches(12.0), Inches(4.2))
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run(); r.text = "•  " + b
        r.font.size = Pt(15); r.font.color.rgb = INK
    # footer
    ftf = _box(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))
    fr = ftf.paragraphs[0].add_run()
    fr.text = "Deniz Ege | Methodology & Simulation | backup"
    fr.font.size = Pt(10); fr.font.color.rgb = GREY


# ─────────────────────────── content (verified numbers) ───────────────────────────
CARDS = [
    ("Is this a trace replay of the logs, or a real model?",
     "Distribution-driven, not a replay.",
     ["I fit distributions (inter-batch gap, batch size, kit composition, pick times) "
      "from 4 months of the ZWM92 dispatch log: 167,784 rows, 40,804 kit orders, 103 active days.",
      "The model SAMPLES from those distributions, it does not play history back.",
      "That is exactly what lets me move materials to new locations and re-run the SAME "
      "warehouse under a different policy. A replay could never test a new layout.",
      "Built in Python with SimPy (discrete-event simulation)."]),

    ("Why discrete-event simulation instead of a formula or a real pilot?",
     "The bottleneck is queueing, which only DES captures.",
     ["The delay is operators waiting for a shared reach-truck fleet, plus a cascade: "
      "an order holds its slot while it waits, delaying the next order.",
      "Static/analytic models assume average flow and miss this contention.",
      "A real pilot means physically moving pallets first; the simulation tests the "
      "policy at zero cost and zero risk, then we recommend a pilot as future work.",
      "Shared resources modelled: 8 operators, 7 reach trucks, 4 Kardex units, pallet-position locks."]),

    ("Where do the pick / handling times come from?",
     "From a video time-motion study of the F400 line.",
     ["2,319 micro-events, 296.6 observed minutes, coded from the F400 kitting video.",
      "Service times are lognormal with spread set by the measured CV:  sigma = sqrt(ln(1 + CV^2)).",
      "Means: operator pick 0.113 min, reach-truck pick+place 0.110, manual penalty 0.102, Kardex pick 0.113.",
      "Sigmas: operator 1.245, reach-truck 1.047, manual 1.279 (moment-matched from the video CVs)."]),

    ("How did you validate the Kardex timing? There was no Kardex video.",
     "Honest answer: it is an assumption, and it cannot bias the comparison.",
     ["No Kardex-specific footage existed, so Kardex pick mean and CV MIRROR the F400 "
      "operator pick (0.113 min, sigma 1.245); carousel turn 0.4 min, sigma 0.30 (book value).",
      "This is disclosed openly in the report and assumptions.",
      "Kardex routing is identical across all six policies, so any timing error cancels "
      "out of the policy comparison; it shifts levels, not the ranking."]),

    ("F400 timing is extrapolated to 9 product families. What if F400 is unrepresentative?",
     "Sensitivity analysis bounds the risk; rankings are preserved.",
     ["F400 is about 29% of dispatch rows; its timing is applied to the other 8 families.",
      "A one-at-a-time sensitivity sweep shifts every timing input by +/-20%.",
      "Line-aware Slotting stays the #1 policy under every swing; lead time moves "
      "proportionally, the conclusion does not flip.",
      "Per-line video time studies are listed as future work to retire this limitation."]),

    ("Why model batch arrivals, and does the arrival rate match reality?",
     "Because the real log arrives in batches, and the rate calibrates within ~1%.",
     ["ZWM92 shows 3,885 same-timestamp batches, mean 4.68 kits per batch.",
      "Inter-batch gap mean 5.36 min with CV 2.04, so it is empirical, not Exponential.",
      "Analytic rate is about 419 orders/day; the full 20-replication runs measure "
      "~398-400/day, within 1.1% of the real 396/day average.",
      "A single short 2-day check gave 378/day (-4.6%); all are inside the +/-10% band."]),

    ("Why 20 replications, and why the same seeds across policies?",
     "Common Random Numbers, and 20 runs separate the policies cleanly.",
     ["Each policy runs N=20 replications; replication i uses seed 42+i, with separate "
      "streams for arrivals and service times.",
      "CRN means every policy faces the EXACT same demand, so differences are the policy, not luck.",
      "6 policies x 20 = 120 runs (Minitab-ready).",
      "95% confidence intervals (t, df=19) for Line-aware and the baselines are narrow and "
      "do NOT overlap, which is why 20 is enough."]),

    ("Your chi-square is highly significant. Doesn't that mean the model is wrong?",
     "No. Significance is expected here; read the effect size.",
     ["Pooled chi^2(7) = 535.05, p < 0.001, but Cramer's V = 0.22 (small-to-moderate effect).",
      "It is significant because the model uses its OWN routing logic, not a bin-for-bin replay; "
      "an exact match would only happen if we replayed history, defeating the purpose.",
      "This is an internal-consistency check, not out-of-sample validation.",
      "Other checks: daily volume within 1.1%, and zero rack crossings in the 3D route audit."]),

    ("You fitted and validated on the same ZWM92 data. No holdout?",
     "Correct, and it is disclosed; no independent period was available.",
     ["SAP never logged waiting times or reach-truck calls, so there was no separate "
      "observation period to hold out.",
      "Independent grounding still exists: the F400 video anchors the timing, daily volume "
      "calibrates to within 1.1%, and the 3D viewer confirms physically valid routes.",
      "We label these as internal-consistency tests, not a true holdout, on the limitations slide."]),

    ("Aisle congestion is not modelled. Doesn't that invalidate the results?",
     "It makes absolute times slightly optimistic, but the ranking is safe.",
     ["Operators and trucks pass freely in shared corridors, so high-load lead times are "
      "a little optimistic.",
      "All six policies share this assumption symmetrically, so the comparison is fair.",
      "Adding congestion would only WIDEN Line-aware's lead, since the baselines saturate "
      "the reach-truck fleet far more (Travel-distance hits 44% RT utilization)."]),

    ("How real is the 'Actual SAP' baseline?",
     "It is an honestly-labelled SAP-plus-heuristic hybrid.",
     ["750 materials sit at their exact decoded SAP rack-bay-position; 2,872 are Kardex-routed "
      "(policy-invariant); the rest fall to an FMR heuristic fallback.",
      "386 materials have multiple SAP bins and are not split into forward vs reserve.",
      "We never claim a pure bin-for-bin replay; it is the best reconstruction the decoded "
      "SAP data allows, and it is the validation ground truth."]),

    ("Throughput barely changes and RT utilization is only 1.3%. Is that believable?",
     "Yes: the system is demand-limited, and 1.3% is picking-only.",
     ["Every policy sees the same arrival stream, so throughput is ~398/day for all of them "
      "(ANOVA F about 0.01, p about 1.0). The win is in lead/wait, not output.",
      "RT utilization 1.3% reflects PICKING workload only; replenishment is excluded for "
      "ALL policies equally, so the comparison stays fair.",
      "In a live deployment total RT load would be higher once replenishment is added; this "
      "is the key caveat we disclose."]),

    ("What about the milk-run trains?",
     "Defined in the model, but operator-carried in the headline runs.",
     ["The model represents 7 milk-run trains as an available transport mode.",
      "In the reported configuration the operator carries the finished kit to the line "
      "(milk-run delivery is switched off), so reported lead time is attributed to operator "
      "and reach-truck only.",
      "This is a modelling choice, disclosed; enabling milk-run would offload operator "
      "walking, it does not change the reach-truck queue that drives the result."]),

    ("Is the warehouse layout real or invented?",
     "Reconstructed from the real AutoCAD drawing and rack PDFs.",
     ["Geometry from the AutoCAD DWG plus 11 per-rack pallet PDFs, cross-checked by a blind "
      "CAD extraction.",
      "11 rack rows (A-J and U); 3,101 modelled pallet positions reconcile with the 3,203 "
      "drawing capacity (the 102 gap = narrower feeder bays plus a cart-parking slot).",
      "5,941 active materials processed; this is a faithful digital twin of the floor."]),
]

SUMMARY = [
    ("Lead time", "6.91 min   [6.44, 7.38]", "12.37 min", "-44%"),
    ("Operator waiting", "4.26 min", "8.45 min", "-50%"),
    ("Reach-truck utilization", "1.27%", "21.71%", "-20.4 pts"),
    ("Operator utilization", "27.5%", "40.5%", "lower"),
    ("Walking distance", "92.4 m", "89.0 m", "+3.8%"),
    ("Throughput", "398.6 / day", "398.5 / day", "unchanged"),
]


def summary_slide():
    s = prs.slides.add_slide(BLANK)
    band = _rect(s, 0, 0, SW, Inches(1.2), NAVY)
    band.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    band.text_frame.margin_left = Inches(0.6)
    r = band.text_frame.paragraphs[0].add_run()
    r.text = "One-glance numbers (Line-aware vs SAP baseline, N=20 CRN)"
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    rows = len(SUMMARY) + 1
    tbl = s.shapes.add_table(rows, 4, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.4)).table
    heads = ["KPI", "Line-aware", "SAP baseline", "Change"]
    for c, h in enumerate(heads):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5E)
        pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.color.rgb = WHITE
        pr.runs[0].font.bold = True; pr.runs[0].font.size = Pt(14)
    for ri, row in enumerate(SUMMARY, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(13)
            run.font.color.rgb = GREEN if ci == 1 else INK
            run.font.bold = ci in (0, 1)
    ftf = _box(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.9))
    for i, line in enumerate([
        "Effect size: Cohen's d = 3.10. Paired-t vs SAP: t about -17.5, p about 3.7e-13.",
        "ANOVA: lead F=47.07 (eta^2=0.67), RT util F=430 (eta^2=0.95), throughput F about 0.01 (no effect).",
    ]):
        p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.size = Pt(12); run.font.color.rgb = GREY


# ─────────────────────────── build ───────────────────────────
title_slide()
for i, (q, a, bs) in enumerate(CARDS, start=1):
    qna_slide(i, q, a, bs)
summary_slide()

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
